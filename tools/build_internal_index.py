import base64
import csv
import hashlib
import json
import os
import re
from datetime import datetime
from pathlib import Path

from cryptography.hazmat.primitives.ciphers.aead import AESGCM

try:
    import openpyxl
except Exception:
    openpyxl = None

try:
    from docx import Document
except Exception:
    Document = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None

try:
    import pdfplumber
except Exception:
    pdfplumber = None


WORKSPACE = Path(__file__).resolve().parents[1]
SOURCE_ROOT = Path(os.environ.get("WINDER_SOURCE_DIR", WORKSPACE.parent / "00_分类整理")).resolve()
DATA_DIR = WORKSPACE / "data"
PASSWORD = "Winder-KB-2026-7F4M-92QX"
CREATED_AT = datetime.now().strftime("%Y-%m-%d %H:%M:%S")

CATEGORY_RULES = [
    ("PLC学习资料", {
        "欧姆龙 NX701-1600": ["nx701", "nj", "omron", "欧姆龙"],
        "Sysmac Studio": ["sysmac"],
        "I/O输入输出": ["i/o", "io", "输入", "输出"],
        "EtherCAT": ["ethercat"],
        "ST语言": ["st语言", "structured text"],
        "功能块FB": ["功能块", "fb", "function block"],
        "传感器接线": ["传感器", "接线", "npn", "pnp", "com端"],
    }),
    ("伺服驱动器资料", {
        "汇川伺服": ["汇川", "sv660", "sv630", "inovance"],
        "欧姆龙伺服": ["欧姆龙伺服", "omron servo"],
        "安川伺服": ["安川", "yaskawa"],
        "伺服报警代码": ["伺服报警", "servo alarm"],
        "伺服参数": ["伺服参数", "刚性", "惯量", "增益"],
        "原点回归": ["原点", "回零", "home"],
        "编码器报警": ["编码器", "encoder"],
    }),
    ("卷绕机机构资料", {
        "放卷机构": ["放卷", "换卷"],
        "张力机构": ["张力", "tension"],
        "纠偏机构": ["纠偏", "epc"],
        "卷针机构": ["卷针"],
        "夹针机构": ["夹针"],
        "贴胶机构": ["贴胶", "胶纸", "胶带", "收尾胶"],
        "切刀机构": ["切刀", "刀"],
        "下料机构": ["下料", "出料", "夹爪"],
    }),
    ("CCD视觉资料", {
        "下料CCD": ["下料ccd", "下料"],
        "卷绕CCD": ["卷绕ccd"],
        "线扫CCD": ["线扫", "line scan"],
        "极耳检测": ["极耳", "tab"],
        "Mark孔检测": ["mark", "孔"],
        "Overhang检测": ["overhang"],
        "CCD误判案例": ["误判", "ng", "ccd"],
        "CCD校验SOP": ["ccd校验", "标定", "校准", "cali"],
    }),
    ("气动与真空资料", {
        "气缸": ["气缸"],
        "电磁阀": ["电磁阀"],
        "真空发生器": ["真空发生器"],
        "真空压力开关": ["真空表", "真空压力", "zse"],
        "调压阀": ["调压阀"],
        "节流阀": ["节流阀"],
        "气缸不到位故障": ["不到位", "气缸"],
    }),
    ("报警代码资料", {
        "PLC报警": ["plc报警"],
        "HMI报警": ["hmi", "触摸屏"],
        "伺服报警": ["伺服报警"],
        "CCD报警": ["ccd报警"],
        "MES报警": ["mes"],
        "气压报警": ["气压"],
        "真空报警": ["真空报警"],
        "安全门报警": ["安全门"],
    }),
    ("设备故障案例", {
        "故障现象": ["故障", "问题", "异常"],
        "发生工位": ["工位"],
        "报警信息": ["报警"],
        "可能原因": ["原因"],
        "实际原因": ["实际原因"],
        "处理过程": ["处理", "维修"],
        "预防措施": ["预防", "改善"],
    }),
    ("点检保养资料", {
        "日点检": ["日点检", "一级保养"],
        "周点检": ["周点检", "周清洁"],
        "月保养": ["月保养", "三级保养"],
        "CCD校准": ["ccd校准", "ccd校验", "标定"],
        "张力校准": ["张力校准"],
        "卷针水平度校准": ["卷针校验", "水平"],
        "导辊清洁": ["导辊", "过辊", "清洁"],
    }),
    ("BOM备件资料", {
        "气缸型号": ["气缸型号"],
        "传感器型号": ["传感器型号"],
        "电磁阀型号": ["电磁阀型号"],
        "伺服电机型号": ["伺服电机"],
        "驱动器型号": ["驱动器"],
        "相机光源型号": ["相机", "光源", "镜头"],
        "易损件清单": ["备件", "易损件", "bom", "清单"],
    }),
    ("生产良率资料", {
        "班产量统计": ["班产量", "产量"],
        "良率计算": ["良率"],
        "不良类型占比": ["不良"],
        "OEE": ["oee"],
        "MTBF": ["mtbf"],
        "MTTR": ["mttr"],
        "停机时间统计": ["停机"],
    }),
]

SENSITIVE_PATTERNS = [
    "瑞浦", "REPT", "印尼", "机场", "客户", "供应商", "图纸", "BOM", "备件",
    "FAT", "SAT", "GRR", "CPK", "程序", "备份", "工艺", "参数", "行政",
    "报销", "身份证", "护照", "请假", "转正", "设备编号", "项目", "协议",
    "DWG", "SLDPRT", "STEP", "专案", "CCSQD"
]

TYPE_MAP = {
    ".pdf": "PDF",
    ".doc": "Word",
    ".docx": "Word",
    ".xls": "Excel",
    ".xlsx": "Excel",
    ".csv": "CSV",
    ".ppt": "PPT",
    ".pptx": "PPT",
    ".jpg": "图片",
    ".jpeg": "图片",
    ".png": "图片",
    ".md": "Markdown",
    ".txt": "TXT",
    ".zip": "压缩包",
    ".raw": "视觉原始数据",
    ".data": "设备数据",
    ".dwg": "图纸",
    ".step": "模型",
    ".sldprt": "模型",
}


def normalize(value):
    return str(value or "").lower()


def keywords_from_text(text):
    words = []
    candidates = [
        "PLC", "NX701", "Sysmac", "EtherCAT", "ST", "FB", "传感器", "伺服", "报警",
        "CCD", "线扫", "极耳", "Mark", "Overhang", "气缸", "电磁阀", "真空",
        "张力", "纠偏", "卷针", "贴胶", "切刀", "下料", "SOP", "点检", "保养",
        "BOM", "备件", "OEE", "MTBF", "MTTR", "FAT", "CPK", "GRR"
    ]
    lower = normalize(text)
    for word in candidates:
        if normalize(word) in lower and word not in words:
            words.append(word)
    return words[:12]


def classify(search_text):
    lower = normalize(search_text)
    best = ("其他资料", "待人工确认", 0)
    for category, subrules in CATEGORY_RULES:
        for sub_category, terms in subrules.items():
            score = sum(1 for term in terms if normalize(term) in lower)
            if score > best[2]:
                best = (category, sub_category, score)
    return best[0], best[1]


def file_type(path):
    return TYPE_MAP.get(path.suffix.lower(), path.suffix.lower().lstrip(".") or "未知")


def internal_flags(text, path):
    joined = f"{text} {path.suffix}".lower()
    flags = []
    for pattern in SENSITIVE_PATTERNS:
        if pattern.lower() in joined:
            flags.append(pattern)
    return sorted(set(flags))


def safe_summary(title, category, sub_category, file_kind, flags):
    if flags:
        return f"内部资料索引：{title}。已按“{category}/{sub_category}”归类；公网不提供原文。"
    return f"{file_kind}资料：{title}。已按“{category}/{sub_category}”归类。"


def extract_preview(path):
    suffix = path.suffix.lower()
    try:
        if suffix in [".txt", ".md", ".csv"]:
            raw = path.read_text(encoding="utf-8", errors="ignore")
            return raw[:3000].strip()
        if suffix == ".pdf" and pdfplumber:
            parts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages[:2]:
                    parts.append(page.extract_text() or "")
            return "\n".join(parts)[:3000].strip()
        if suffix == ".docx" and Document:
            doc = Document(path)
            return "\n".join(p.text for p in doc.paragraphs[:60])[:3000].strip()
        if suffix == ".xlsx" and openpyxl:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets[:2]:
                rows.append(f"[{ws.title}]")
                for row in ws.iter_rows(max_row=20, values_only=True):
                    values = [str(v) for v in row if v is not None]
                    if values:
                        rows.append(" | ".join(values[:10]))
            wb.close()
            return "\n".join(rows)[:3000].strip()
        if suffix == ".pptx" and Presentation:
            prs = Presentation(path)
            chunks = []
            for slide in prs.slides[:8]:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return "\n".join(chunks)[:3000].strip()
    except Exception as exc:
        return f"内容提取失败：{type(exc).__name__}"
    return ""


def make_card(index, path):
    rel_parts = path.relative_to(SOURCE_ROOT).parts
    title = path.stem
    ext = path.suffix.lower()
    file_kind = file_type(path)
    search_text = " ".join(rel_parts)
    preview = extract_preview(path)
    combined = f"{search_text} {preview}"
    category, sub_category = classify(combined)
    flags = internal_flags(combined, path)
    # 该来源文件夹来自个人工作资料，默认按内部资料处理；公开站点不暴露原件。
    is_internal = True
    keywords = keywords_from_text(combined)
    if sub_category not in keywords:
        keywords.insert(0, sub_category)
    if category not in keywords:
        keywords.insert(0, category)
    doc_id = f"D{index:04d}"
    return {
        "id": doc_id,
        "title": title,
        "category": category,
        "sub_category": sub_category,
        "keywords": keywords[:14],
        "summary": safe_summary(title, category, sub_category, file_kind, flags),
        "source_file": path.name,
        "file_type": file_kind,
        "file_path": f"internal://document/{doc_id}",
        "is_internal": is_internal,
        "created_at": CREATED_AT,
        "size_kb": round(path.stat().st_size / 1024, 1),
        "preview_text": preview[:1200],
        "extract_status": "已提取摘要" if preview and not preview.startswith("内容提取失败") else "原文查看",
        "original_access": "内部原始文件未上传公网",
        "sensitive_flags": flags,
    }


def encrypt_payload(payload):
    data = json.dumps(payload, ensure_ascii=False, separators=(",", ":")).encode("utf-8")
    salt = os.urandom(16)
    iv = os.urandom(12)
    key = hashlib.pbkdf2_hmac("sha256", PASSWORD.encode("utf-8"), salt, 150000, dklen=32)
    cipher = AESGCM(key).encrypt(iv, data, None)
    return {
        "encrypted": True,
        "version": 1,
        "algorithm": "PBKDF2-SHA256-AES-256-GCM",
        "kdf": {
            "name": "PBKDF2",
            "hash": "SHA-256",
            "iterations": 150000,
            "salt": base64.b64encode(salt).decode("ascii"),
        },
        "iv": base64.b64encode(iv).decode("ascii"),
        "ciphertext": base64.b64encode(cipher).decode("ascii"),
    }


def write_encrypted(name, payload):
    (DATA_DIR / name).write_text(
        json.dumps(encrypt_payload(payload), ensure_ascii=False, indent=2),
        encoding="utf-8",
    )


def alarm_rows(cards):
    rows = []
    for card in cards:
        text = f"{card['title']} {' '.join(card['keywords'])}"
        if "报警" not in text.lower() and "alarm" not in text.lower():
            continue
        rows.append({
            "document_id": card["id"],
            "brand": infer_brand(text),
            "model": infer_model(text),
            "alarm_code": "",
            "alarm_name": card["title"],
            "cause": "需登录后结合原始内部资料确认。",
            "check_method": "按资料卡片关键词和原始文件名定位后人工核对。",
            "solution": "内部资料未上传公网，需在受控环境查看原文。",
        })
    return rows


def fault_rows(cards):
    rows = []
    for card in cards:
        text = f"{card['title']} {' '.join(card['keywords'])}"
        if not any(word in text for word in ["故障", "问题", "异常", "FAT"]):
            continue
        rows.append({
            "document_id": card["id"],
            "fault_name": card["title"],
            "station": card["sub_category"],
            "symptom": card["summary"],
            "possible_causes": ["需结合原文确认"],
            "actual_cause": "",
            "solution": "内部资料未上传公网，需在受控环境查看原文。",
            "prevention": "",
        })
    return rows


def infer_brand(text):
    for brand in ["欧姆龙", "汇川", "安川", "Keyence", "SMC", "亚德客"]:
        if brand.lower() in text.lower():
            return brand
    return ""


def infer_model(text):
    models = re.findall(r"(NX701|SV660|SV630|MD500|ZSE20F|[A-Z]{1,4}-?\d{2,5})", text, re.I)
    return models[0] if models else ""


def main():
    DATA_DIR.mkdir(exist_ok=True)
    files = sorted(p for p in SOURCE_ROOT.rglob("*") if p.is_file())
    cards = [make_card(i + 1, path) for i, path in enumerate(files)]

    index_payload = {
        "generated_at": CREATED_AT,
        "source": "00_分类整理",
        "security_note": "该文件夹来源于个人工作资料，索引已加密；公网不上传原始文件和真实本机路径。",
        "total": len(cards),
        "documents": cards,
    }

    by_category = {name: [] for name, _ in CATEGORY_RULES}
    for card in cards:
        by_category.setdefault(card["category"], []).append(card)

    write_encrypted("document_index.json", index_payload)
    write_encrypted("alarm_codes.json", alarm_rows(cards))
    write_encrypted("fault_cases.json", fault_rows(cards))
    write_encrypted("plc_notes.json", by_category.get("PLC学习资料", []))
    write_encrypted("servo_notes.json", by_category.get("伺服驱动器资料", []))
    write_encrypted("ccd_notes.json", by_category.get("CCD视觉资料", []))
    write_encrypted("pneumatic_notes.json", by_category.get("气动与真空资料", []))
    write_encrypted("bom_parts.json", by_category.get("BOM备件资料", []))
    write_encrypted("maintenance.json", by_category.get("点检保养资料", []))
    write_encrypted("production_quality.json", by_category.get("生产良率资料", []))

    report = {
        "generated_at": CREATED_AT,
        "total_files": len(cards),
        "internal_files": sum(1 for c in cards if c["is_internal"]),
        "sensitive_files": sum(1 for c in cards if c["sensitive_flags"]),
        "category_counts": {k: len(v) for k, v in sorted(by_category.items())},
        "security_decision": [
            "不上传 00_分类整理 原始文件。",
            "不写入真实本机路径。",
            "内部资料索引写入加密 JSON，登录后才在浏览器解密。",
            "GitHub Pages 前端登录仍不是服务器级权限控制。",
        ],
    }
    (DATA_DIR / "security_report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"generated {len(cards)} encrypted document cards")


if __name__ == "__main__":
    main()
